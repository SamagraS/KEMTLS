from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Image paths ──
IMG_DIR = r"C:\Users\Samagra\.gemini\antigravity\brain\5afb85a0-5d7a-4def-98f5-5e2b3cb8e62c"
IMAGES = {
    1: None,  # will find dynamically
    2: None,
    3: None,
    4: None,
    5: None,
    6: None,
    7: None,
    8: None,
    9: None,
    10: None,
}

# Find images by prefix
for f in os.listdir(IMG_DIR):
    if f.endswith('.png'):
        if 'slide1_hero_' in f and '1775613638' in f: IMAGES[1] = os.path.join(IMG_DIR, f)
        elif 'slide2_cracked_lock_' in f: IMAGES[2] = os.path.join(IMG_DIR, f)
        elif 'slide3_stack_' in f: IMAGES[3] = os.path.join(IMG_DIR, f)
        elif 'slide4_flow_' in f: IMAGES[4] = os.path.join(IMG_DIR, f)
        elif 'slide5_arch_' in f: IMAGES[5] = os.path.join(IMG_DIR, f)
        elif 'slide6_novelties_' in f: IMAGES[6] = os.path.join(IMG_DIR, f)
        elif 'slide7_feasibility_' in f: IMAGES[7] = os.path.join(IMG_DIR, f)
        elif 'slide8_benchmark_' in f: IMAGES[8] = os.path.join(IMG_DIR, f)
        elif 'slide9_impact_' in f: IMAGES[9] = os.path.join(IMG_DIR, f)
        elif 'slide10_closing_' in f: IMAGES[10] = os.path.join(IMG_DIR, f)

# ── Colors ──
PRIMARY_BLUE = RGBColor(0x1A, 0x56, 0xDB)      # #1A56DB
DARK_TEXT = RGBColor(0x1F, 0x2A, 0x37)            # #1F2A37
BODY_TEXT = RGBColor(0x37, 0x41, 0x51)            # #374151
SUBTITLE_TEXT = RGBColor(0x6B, 0x72, 0x80)        # #6B7280
ACCENT_TEAL = RGBColor(0x06, 0x94, 0xA2)         # #0694A2
ACCENT_INDIGO = RGBColor(0x56, 0x50, 0xD6)       # #5650D6
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF9, 0xFA, 0xFB)            # #F9FAFB
ACCENT_BAR = RGBColor(0x1A, 0x56, 0xDB)          # blue accent bar
BULLET_BLUE = RGBColor(0x2C, 0x5F, 0xE0)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


def add_accent_bar(slide, top=Inches(0), left=Inches(0), width=Inches(0.08), height=Inches(7.5), color=ACCENT_BAR):
    """Add a thin colored accent bar on the left edge."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_bottom_bar(slide, color=ACCENT_BAR):
    """Add a thin colored bar at the bottom."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.2), SLIDE_WIDTH, Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_slide_number(slide, num):
    """Add slide number in bottom-right."""
    txBox = slide.shapes.add_textbox(Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(11)
    p.font.color.rgb = SUBTITLE_TEXT
    p.alignment = PP_ALIGN.RIGHT


def set_text_props(run, size=Pt(16), color=BODY_TEXT, bold=False, font_name='Calibri'):
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name


def add_title_textbox(slide, text, left=Inches(0.6), top=Inches(0.5), width=Inches(7), height=Inches(0.8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_text_props(run, size=Pt(30), color=DARK_TEXT, bold=True, font_name='Calibri')
    return txBox


def add_subtitle_textbox(slide, text, left=Inches(0.6), top=Inches(1.3), width=Inches(7), height=Inches(0.5)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_text_props(run, size=Pt(14), color=SUBTITLE_TEXT, bold=False, font_name='Calibri')
    return txBox


def add_bullets(slide, bullets, left=Inches(0.8), top=Inches(2.0), width=Inches(7.0), height=Inches(4.5),
                font_size=Pt(15), line_spacing=Pt(28), bold_prefix=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.space_after = Pt(6)
        p.space_before = Pt(2)

        # Check if bullet has bold prefix like "**KEMTLS-PDK:**"
        if bold_prefix and ':' in bullet:
            parts = bullet.split(':', 1)
            run_bold = p.add_run()
            run_bold.text = "• " + parts[0].strip() + ": "
            set_text_props(run_bold, size=font_size, color=PRIMARY_BLUE, bold=True)
            run_rest = p.add_run()
            run_rest.text = parts[1].strip()
            set_text_props(run_rest, size=font_size, color=BODY_TEXT, bold=False)
        else:
            run = p.add_run()
            run.text = "•  " + bullet
            set_text_props(run, size=font_size, color=BODY_TEXT, bold=False)

    return txBox


def add_image(slide, img_path, left=Inches(8.2), top=Inches(1.5), width=Inches(4.5)):
    if img_path and os.path.exists(img_path):
        slide.shapes.add_picture(img_path, left, top, width=width)


# ══════════════════════════════════════════════
# SLIDE 1: Title Slide
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_accent_bar(slide)
add_bottom_bar(slide)

# Title
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(7.5), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Post-Quantum OpenID Connect\nover KEMTLS and QUIC"
set_text_props(run, size=Pt(36), color=DARK_TEXT, bold=True)

# Subtitle
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(7.5), Inches(1.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Session-Bound Tokens, KEMTLS-PDK, and Research-Grade Benchmarking"
set_text_props(run, size=Pt(18), color=SUBTITLE_TEXT, bold=False)

# Team / submission
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(5), Inches(1.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Team Name  |  Submission / Challenge Name"
set_text_props(run, size=Pt(14), color=SUBTITLE_TEXT, bold=False)

# Image
add_image(slide, IMAGES[1], left=Inches(8.5), top=Inches(1.2), width=Inches(4.2))
add_slide_number(slide, 1)


# ══════════════════════════════════════════════
# SLIDE 2: Problem and Opportunity
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_accent_bar(slide)
add_bottom_bar(slide)
add_title_textbox(slide, "Problem and Opportunity")
add_subtitle_textbox(slide, "Why current identity systems need a post-quantum transport redesign")

add_bullets(slide, [
    "Modern identity systems still rely on classical public-key cryptography.",
    "Post-quantum migration often keeps the traditional secure transport model unchanged.",
    "KEMTLS is a strong post-quantum transport candidate, but is underexplored in OIDC.",
    "Opportunity: preserve OpenID Connect while improving transport security and token protection."
])

add_image(slide, IMAGES[2], left=Inches(8.2), top=Inches(1.5), width=Inches(4.5))
add_slide_number(slide, 2)


# ══════════════════════════════════════════════
# SLIDE 3: Proposed Solution
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_accent_bar(slide)
add_bottom_bar(slide)
add_title_textbox(slide, "Proposed Solution")
add_subtitle_textbox(slide, "Standard OIDC flow, redesigned transport and cryptographic foundation")

add_bullets(slide, [
    "Preserve the standard OIDC Authorization Code Flow with PKCE.",
    "Secure communication using KEMTLS, with QUIC integrated into the transport path.",
    "ML-KEM-768 for key establishment and ML-DSA-65 for token signatures.",
    "Python for OIDC orchestration + Rust for high-performance transport.",
    "Strengthen token security without changing the visible OpenID Connect flow."
])

add_image(slide, IMAGES[3], left=Inches(8.2), top=Inches(1.5), width=Inches(4.5))
add_slide_number(slide, 3)


# ══════════════════════════════════════════════
# SLIDE 4: End-to-End System Flow
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_accent_bar(slide)
add_bottom_bar(slide)
add_title_textbox(slide, "End-to-End System Flow")
add_subtitle_textbox(slide, "From login to protected resource access and token refresh")

add_bullets(slide, [
    "Client establishes a secure session with the authorization server.",
    "Authorization and token exchange run through the protected transport channel.",
    "Auth server issues post-quantum-signed identity and access tokens.",
    "Client accesses the resource server through the same secure model.",
    "Resource server verifies both token validity and live session binding.",
    "Refresh-token rotation rebinds newly issued tokens to the active session."
])

add_image(slide, IMAGES[4], left=Inches(8.2), top=Inches(1.5), width=Inches(4.5))
add_slide_number(slide, 4)


# ══════════════════════════════════════════════
# SLIDE 5: Technical Architecture
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_accent_bar(slide)
add_bottom_bar(slide)
add_title_textbox(slide, "Technical Architecture")
add_subtitle_textbox(slide, "Modular Python + Rust implementation stack")

add_bullets(slide, [
    "Python handles OIDC orchestration, endpoint logic, demo flow, and benchmark control.",
    "Rust handles performance-critical transport and secure communication components.",
    "Modular architecture across cryptography, transport, OIDC, servers, client, and benchmarking.",
    "Supports baseline trust, KEMTLS-PDK, and QUIC-integrated deployment modes."
])

add_image(slide, IMAGES[5], left=Inches(8.2), top=Inches(1.5), width=Inches(4.5))
add_slide_number(slide, 5)


# ══════════════════════════════════════════════
# SLIDE 6: Core Novelties
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_accent_bar(slide)
add_bottom_bar(slide)
add_title_textbox(slide, "Core Novelties")
add_subtitle_textbox(slide, "Five key innovations in this system")

add_bullets(slide, [
    "KEMTLS-PDK: reduces certificate-related trust-establishment overhead in managed environments.",
    "QUIC-Integrated Transport: improves communication efficiency and responsiveness.",
    "Session-Bound Access Tokens: valid tokens fail when replayed on a different session.",
    "KEMTLS-Bound Refresh Rotation: refresh tokens rotate and rebind to the active session.",
    "Research-Grade Benchmarking: evaluates crypto, transport, and full identity-flow behavior."
], bold_prefix=True)

add_image(slide, IMAGES[6], left=Inches(8.2), top=Inches(1.5), width=Inches(4.5))
add_slide_number(slide, 6)


# ══════════════════════════════════════════════
# SLIDE 7: Feasibility and Productization
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_accent_bar(slide)
add_bottom_bar(slide)
add_title_textbox(slide, "Feasibility and Productization")
add_subtitle_textbox(slide, "Deployment paths, challenges, and mitigations")

add_bullets(slide, [
    "Suitable for enterprise identity, regulated environments, and controlled PQ migration.",
    "Supports certificate-oriented trust and tightly managed KEMTLS-PDK deployments.",
    "QUIC integration provides a practical path for performance-oriented secure systems.",
    "Challenges: browser compatibility, interoperability, deployment complexity.",
    "Mitigations: controlled pilots, modular integration, standards-aligned flow, phased adoption."
])

add_image(slide, IMAGES[7], left=Inches(8.2), top=Inches(1.5), width=Inches(4.5))
add_slide_number(slide, 7)


# ══════════════════════════════════════════════
# SLIDE 8: Benchmarking and Research Methodology
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_accent_bar(slide)
add_bottom_bar(slide)
add_title_textbox(slide, "Benchmarking and Research Methodology")
add_subtitle_textbox(slide, "A layered, reproducible evaluation framework")

add_bullets(slide, [
    "Multi-Layer Evaluation: crypto primitives → transport → full OIDC flow.",
    "Comparison Model: local results vs. classical and PQ literature baselines.",
    "Measurement Discipline: warm-up exclusion, repeated trials, per-phase timing, size analysis.",
    "Reproducibility: Python+Rust modular stack, fixed config, raw and processed outputs preserved."
], bold_prefix=True)

add_image(slide, IMAGES[8], left=Inches(8.2), top=Inches(1.5), width=Inches(4.5))
add_slide_number(slide, 8)


# ══════════════════════════════════════════════
# SLIDE 9: Performance, Impact, and Beneficiaries
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_accent_bar(slide)
add_bottom_bar(slide)
add_title_textbox(slide, "Performance, Impact, and Beneficiaries")
add_subtitle_textbox(slide, "Measurable improvements and target sectors")

add_bullets(slide, [
    "KEMTLS-PDK reduces trust-establishment overhead in managed deployments.",
    "QUIC-integrated transport improves efficiency and responsiveness.",
    "Session-bound tokens reduce replay risk and strengthen resource access.",
    "Practical migration path toward post-quantum-secure identity infrastructure.",
    "Key sectors: banking, government, healthcare, and enterprise identity."
])

add_image(slide, IMAGES[9], left=Inches(8.2), top=Inches(1.5), width=Inches(4.5))
add_slide_number(slide, 9)


# ══════════════════════════════════════════════
# SLIDE 10: Research References and Closing
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_accent_bar(slide)
add_bottom_bar(slide)
add_title_textbox(slide, "Research References and Closing")
add_subtitle_textbox(slide, "Building on KEMTLS research, NIST standards, and post-quantum OIDC studies")

add_bullets(slide, [
    "Builds on KEMTLS research, PQ-OIDC studies, and NIST standards for ML-KEM / ML-DSA.",
    "Preserves OIDC semantics while redesigning transport and token security for the PQ era.",
    "Main contribution: KEMTLS + KEMTLS-PDK + QUIC + session-bound tokens + benchmarking.",
    "[Insert Repository / Demo Link / QR Code Here]"
])

# Closing statement
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(8), Inches(1))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Thank you."
set_text_props(run, size=Pt(28), color=PRIMARY_BLUE, bold=True)

add_image(slide, IMAGES[10], left=Inches(9.0), top=Inches(3.5), width=Inches(3.5))
add_slide_number(slide, 10)


# ── Save ──
output_path = r"d:\KEMTLS\Post_Quantum_OIDC_Presentation.pptx"
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"   Slides: {len(prs.slides)}")
for i, img in IMAGES.items():
    status = "✓" if img and os.path.exists(img) else "✗ MISSING"
    print(f"   Slide {i} image: {status}")
